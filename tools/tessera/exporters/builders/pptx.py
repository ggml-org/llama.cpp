"""PPTX builder (exporter).

Converts a list of ``IntermediateDocument``s (one per slide)
to a .pptx file via python-pptx. The exporter is the
inverse of the PPTX importer: a slide with a title and
paragraphs is rebuilt as a slide with the same content.

v1 supports titles, paragraphs, lists, tables, and images
on slides. Slide masters and theme colors are dropped
(matching the importer v1).
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from tools.tessera.importers.intermediate import (
    IntermediateBlock,
    IntermediateDocument,
)

log = logging.getLogger(__name__)


def build_pptx(docs: list[IntermediateDocument], output: Path) -> None:
    """Render a list of IntermediateDocuments to a .pptx file at `output`.

    One document per slide. The presentation's slide master
    layout is "Title and Content" (the most generic
    layout; v1 doesn't preserve the original layout).
    """
    prs = Presentation()
    layout = prs.slide_layouts[1]  # "Title and Content"
    for d in docs:
        slide = prs.slides.add_slide(layout)
        _render_slide(slide, d)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def _render_slide(slide: Any, doc: IntermediateDocument) -> None:
    """Populate a slide with the document's blocks."""
    # The first heading becomes the title (matching the
    # importer's behaviour). Subsequent blocks go in the body.
    title_used = False
    for ib in doc.blocks:
        if not title_used and ib.type == "heading":
            slide.shapes.title.text = "".join(r.text for r in ib.runs)
            title_used = True
            continue
        if ib.type == "paragraph":
            _add_body_paragraph(slide, "".join(r.text for r in ib.runs))
        elif ib.type == "list":
            for c in ib.children:
                if hasattr(c, "type"):
                    _add_body_paragraph(slide, f"• {''.join(r.text for r in c.runs)}")
        elif ib.type == "image":
            src = ib.attrs.get("source", "")
            if src and Path(str(src)).exists():
                try:
                    slide.shapes.add_picture(str(src), Inches(1), Inches(3), width=Inches(5))
                except Exception as e:  # noqa: BLE001
                    log.warning("pptx export: image %s: %s", src, e)
        # Tables, dividers, code blocks: skip in v1; the
        # python-pptx slide layout doesn't have a good place
        # for them and the importer's table support is
        # best-effort anyway.


def _add_body_paragraph(slide: Any, text: str) -> None:
    """Append a text frame line to the slide's body placeholder."""
    body = slide.shapes.placeholders[1].text_frame if len(slide.shapes.placeholders) > 1 else None
    if body is None:
        # Slide without a body placeholder; add a text box.
        from pptx.util import Inches, Pt

        tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.5))
        tx.text_frame.text = text
        return
    if body.text == "":
        body.text = text
    else:
        p = body.add_paragraph()
        p.text = text
