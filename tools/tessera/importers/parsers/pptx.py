"""PPTX parser.

Reads a .pptx with python-pptx and produces one
``IntermediateDocument`` per slide. Each slide is a sequence of
blocks: a title heading (when present), then one paragraph or
list per text frame.

The v1 PPTX importer is deliberately simple:

* Each slide becomes one document (the productivity surface's
  "slide" entity type). Slides are linked in a chain via
  ``entity_links`` so the user can navigate the deck in order.
* Text frames become paragraph blocks. Bulleted text frames
  become a single ``list`` block with one ``listItem`` per
  bullet.
* Tables become a single ``table`` block per slide.
* Pictures become an ``image`` block. Image bytes are written
  to the media directory the caller provides (we extract from
  the underlying zip).

Punted (v1):

* Slide master layouts and theme colors.
* Speaker notes (python-pptx exposes them but we don't surface
  them; v2 adds a ``notes`` block type).
* Transitions, animations, embedded video / audio.
"""

from __future__ import annotations

import logging
import shutil
import zipfile
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Emu

from ..intermediate import (
    IntermediateBlock,
    IntermediateDocument,
    IntermediateInlineRun,
)

log = logging.getLogger(__name__)


def parse_pptx(path: Path, *, media_dir: Optional[Path] = None) -> list[IntermediateDocument]:
    """Parse a .pptx file into one IntermediateDocument per slide.

    `media_dir` is where image bytes are extracted. When None,
    pictures are skipped (a warning is logged).
    """
    log.debug("parse_pptx: %s", path)
    prs = Presentation(str(path))
    docs: list[IntermediateDocument] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        docs.append(_slide_to_doc(slide, slide_idx, path, media_dir))
    return docs


def _slide_to_doc(
    slide: Any, slide_idx: int, path: Path, media_dir: Optional[Path]
) -> IntermediateDocument:
    """Convert one python-pptx Slide to an IntermediateDocument.

    We walk the slide shapes in shape-tree order (not the visual
    order) so the output matches the file's structure. The
    ``title`` shape is emitted as a level-1 heading; everything
    else follows.
    """
    blocks: list[IntermediateBlock] = []

    title = slide.shapes.title
    if title is not None and title.has_text_frame:
        for para in title.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs)
            if not text.strip():
                continue
            blocks.append(
                IntermediateBlock(
                    type="heading",
                    attrs={"level": 1},
                    runs=[IntermediateInlineRun(text=text)],
                )
            )
            # The title is a single paragraph; don't keep going
            # with the rest of the title's runs.
            break

    # Group shapes by type so we emit blocks in a sensible order:
    # text frames (lists + paragraphs) first, then tables, then
    # pictures. The visual order is less useful for the AST
    # because the user can re-arrange via the editor.
    text_blocks: list[IntermediateBlock] = []
    table_blocks: list[IntermediateBlock] = []
    image_blocks: list[IntermediateBlock] = []
    for shape in slide.shapes:
        if shape == title:
            continue
        if shape.has_text_frame:
            text_blocks.extend(_text_frame_blocks(shape))
        elif shape.has_table:
            table_blocks.append(_table_to_block(shape))
        elif shape.shape_type == 13:  # PICTURE
            image_blocks.extend(_picture_blocks(shape, slide_idx, path, media_dir))

    blocks.extend(text_blocks)
    blocks.extend(table_blocks)
    blocks.extend(image_blocks)
    return IntermediateDocument(blocks=blocks, meta={"slide_index": slide_idx})


def _text_frame_blocks(shape: Any) -> list[IntermediateBlock]:
    """Convert a text frame into a sequence of blocks.

    A text frame with bullet-style paragraphs becomes a single
    ``list`` block (one item per non-empty paragraph). A text
    frame with plain paragraphs becomes one ``paragraph`` block
    per non-empty paragraph.

    python-pptx's ``paragraph.level`` is the indent level (0 = top
    level). We collapse all levels into a single flat list for v1
    (multi-level lists are punted to v2).
    """
    paragraphs = list(shape.text_frame.paragraphs)
    if not paragraphs:
        return []

    # Decide list vs paragraph: if every paragraph has a
    # non-None bullet, it's a list.
    is_list = all(_has_bullet(p) for p in paragraphs if p.text.strip())
    if is_list:
        items: list[IntermediateBlock] = []
        for p in paragraphs:
            text = p.text.strip()
            if not text:
                continue
            items.append(
                IntermediateBlock(
                    type="listItem",
                    runs=[IntermediateInlineRun(text=text)],
                )
            )
        if not items:
            return []
        return [
            IntermediateBlock(
                type="list",
                attrs={"style": "unordered", "items": [it.id for it in items]},
                children=items,
            )
        ]

    out: list[IntermediateBlock] = []
    for p in paragraphs:
        text = p.text
        if not text.strip():
            continue
        out.append(
            IntermediateBlock(
                type="paragraph",
                runs=[IntermediateInlineRun(text=text)],
            )
        )
    return out


def _has_bullet(p: Any) -> bool:
    """True if a paragraph has a bullet / list marker."""
    # python-pptx's bullet is on paragraph._pPr; a None pPr
    # means no list info.
    pPr = p._pPr
    if pPr is None:
        return False
    return pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar") is not None or \
        pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum") is not None


def _table_to_block(shape: Any) -> IntermediateBlock:
    """Convert a PPTX table shape to a ``table`` block."""
    tbl = shape.table
    rows = list(tbl.rows)
    if not rows:
        return IntermediateBlock(type="table", attrs={"rows": 0, "cols": 0, "cells": []})

    n_rows = len(rows)
    n_cols = len(list(rows[0].cells))
    grid: list[list[IntermediateBlock]] = []
    for r in rows:
        row_blocks: list[IntermediateBlock] = []
        for c in r.cells:
            text = c.text.strip()
            row_blocks.append(
                IntermediateBlock(
                    type="paragraph",
                    runs=[IntermediateInlineRun(text=text)] if text else [],
                )
            )
        # Pad to n_cols in case the table is ragged.
        while len(row_blocks) < n_cols:
            row_blocks.append(IntermediateBlock(type="paragraph", runs=[]))
        grid.append(row_blocks)

    flat: list[IntermediateBlock] = [c for row in grid for c in row]
    return IntermediateBlock(
        type="table",
        attrs={
            "rows": n_rows,
            "cols": n_cols,
            "cells": [[c.id for c in row] for row in grid],
        },
        children=flat,
    )


def _picture_blocks(
    shape: Any, slide_idx: int, path: Path, media_dir: Optional[Path]
) -> list[IntermediateBlock]:
    """Extract a picture's bytes and return an ``image`` block."""
    if media_dir is None:
        log.warning("parse_pptx: image at slide %d skipped (no media_dir)", slide_idx)
        return []
    media_dir.mkdir(parents=True, exist_ok=True)
    image = shape.image
    ext = image.ext or "bin"
    fname = f"slide{slide_idx}_image_{image.filename or 'pic'}.{ext}"
    dst = media_dir / fname
    try:
        with dst.open("wb") as f:
            f.write(image.blob)
    except OSError as e:
        log.warning("image extract failed: %s: %s", path, e)
        return []
    width = int(shape.width / Emu(1)) if shape.width is not None else None
    height = int(shape.height / Emu(1)) if shape.height is not None else None
    return [
        IntermediateBlock(
            type="image",
            attrs={"source": str(dst), "alt": ""},
            meta={"width_emu": width, "height_emu": height},
        )
    ]
